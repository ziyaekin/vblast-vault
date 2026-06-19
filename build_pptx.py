#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
V-BLAST sunumu (.pptx) üretici — python-pptx ile.
23 slayt, Türkçe, formüller (Unicode), makale görselleri + kendi sim grafikleri,
her slaytta anlatım scripti = speaker notes.
"""
import glob, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ART = os.path.join(HERE, 'wolnianskyandfoschini_artifacts')
FIG = os.path.join(HERE, 'sim', 'figures')
F1 = glob.glob(os.path.join(ART, 'image_000000*'))[0]   # Şekil 1 blok diyagram
F2 = glob.glob(os.path.join(ART, 'image_000001*'))[0]   # Şekil 2 BLER-SNR
F3 = glob.glob(os.path.join(ART, 'image_000002*'))[0]   # Şekil 3 konum
SIM1 = os.path.join(FIG, 'SIM-1.png')
SIM2 = os.path.join(FIG, 'SIM-2.png')
SIM3 = os.path.join(FIG, 'SIM-3.png')

# --- renk paleti ---
PRIMARY = RGBColor(0x0B, 0x3D, 0x5C)   # koyu lacivert
ACCENT  = RGBColor(0xC0, 0x39, 0x2B)   # kırmızı vurgu
DARK    = RGBColor(0x22, 0x2A, 0x32)
GRAY    = RGBColor(0x5A, 0x6A, 0x78)
LIGHT   = RGBColor(0xEC, 0xF2, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
FORMBG  = RGBColor(0xF4, 0xF1, 0xE8)   # formül kutusu arka planı

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def set_notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def rect(s, l, t, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def title_band(s, title, num=None):
    rect(s, 0, 0, SW, Inches(1.05), PRIMARY)
    rect(s, 0, Inches(1.05), SW, Inches(0.06), ACCENT)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.0), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    if num is not None:
        nb = s.shapes.add_textbox(Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
        np_ = nb.text_frame.paragraphs[0]; np_.text = str(num)
        np_.font.size = Pt(12); np_.font.color.rgb = GRAY; np_.alignment = PP_ALIGN.RIGHT


def add_bullets(s, l, t, w, h, items, base=20):
    """items: list of (text, level) veya text. level 0/1. '##' başı = vurgulu başlık."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, level = it
        else:
            text, level = it, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        head = text.startswith('##')
        if head:
            text = text[2:].strip()
        bullet = '' if head else ('–  ' if level == 1 else '•  ')
        p.text = bullet + text
        p.level = level
        p.font.size = Pt(base - 2 if level == 1 else base)
        p.space_after = Pt(6)
        if head:
            p.font.bold = True; p.font.color.rgb = ACCENT; p.font.size = Pt(base)
        else:
            p.font.color.rgb = DARK
    return tb


def add_formula(s, l, t, w, lines, size=20):
    """Formülleri belirgin bir kutuda (Unicode metin)."""
    h = Inches(0.55 + 0.42 * len(lines))
    box = rect(s, l, t, w, h, FORMBG)
    box.line.color.rgb = RGBColor(0xD9, 0xCE, 0xB0); box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(8); tf.margin_bottom = Pt(8)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size); p.font.name = 'Consolas'
        p.font.color.rgb = PRIMARY; p.space_after = Pt(4)
    return box


def add_image_fit(s, path, box_l, box_t, box_w, box_h, caption=None):
    """Görseli kutuya en-boy koruyarak sığdır + ortala. caption alt yazı."""
    pic = s.shapes.add_picture(path, box_l, box_t, width=box_w)
    if pic.height > box_h:
        pic.height = box_h; pic.width = int(box_h * (pic.width / pic.height)) if False else pic.width
        # yeniden ölçek: yükseklikten
        ratio = box_h / pic.height
    # en-boy koru: yükseklik taşarsa genişlikten yeniden hesapla
    if pic.height > box_h:
        scale = box_h / pic.height
        pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
    pic.left = int(box_l + (box_w - pic.width) / 2)
    pic.top = int(box_t + (box_h - pic.height) / 2)
    if caption:
        cb = s.shapes.add_textbox(box_l, box_t + box_h, box_w, Inches(0.3))
        cp = cb.text_frame.paragraphs[0]; cp.text = caption
        cp.alignment = PP_ALIGN.CENTER; cp.font.size = Pt(12)
        cp.font.italic = True; cp.font.color.rgb = GRAY
    return pic


# ======================================================================
#  SLAYT 1 — Başlık
# ======================================================================
s = slide()
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, Inches(3.05), SW, Inches(0.08), ACCENT)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'V-BLAST'
p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = 'Zengin-Saçılımlı Kablosuz Kanalda Çok Yüksek Veri Hızları'
p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(0xCF, 0xE0, 0xEC)
sb = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11.5), Inches(2.2))
stf = sb.text_frame; stf.word_wrap = True
for i, line in enumerate([
    'Wolniansky, Foschini, Golden, Valenzuela — Bell Labs, Lucent (1998)',
    '',
    'MIMO • Uzamsal Çoğullama • Nulling + Sıralı İptal',
    'Makale incelemesi + MATLAB simülasyonu',
]):
    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    p.text = line; p.font.size = Pt(18)
    p.font.color.rgb = WHITE if i == 0 else RGBColor(0xBF, 0xD3, 0xE0)
nb = s.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5))
np_ = nb.text_frame.paragraphs[0]; np_.text = 'Sunan: [İsim]  •  Tarih: [….]'
np_.font.size = Pt(16); np_.font.color.rgb = RGBColor(0x9F, 0xB8, 0xC8)
set_notes(s, "Herkese merhaba. Bugün kablosuz haberleşmenin kilometre taşı sayılan bir çalışmayı, "
            "Bell Labs'in V-BLAST mimarisini anlatacağım. Bu makale, aynı frekansı ve aynı zamanı "
            "paylaşan birden fazla anteni kullanarak, kablosuz spektral verimliliği o güne kadar "
            "görülmemiş seviyelere — 20 ila 40 bps/Hz'e — çıkarmayı başardı. Sunumun sonunda, hem "
            "yöntemin nasıl çalıştığını hem de bunu kendi MATLAB simülasyonumuzla nasıl doğruladığımızı "
            "göreceğiz. Önce en temelden, tek antenli sistemlerden başlayacağız.")


# ======================================================================
#  Genel içerik slaytları
# ======================================================================
def content(num, title, bullets=None, formula=None, notes='', base=20,
            formula_top=None):
    s = slide()
    title_band(s, title, num)
    top = Inches(1.35)
    if bullets:
        add_bullets(s, Inches(0.6), top, Inches(12.1), Inches(5.4), bullets, base)
    if formula:
        ft = formula_top if formula_top is not None else Inches(5.5)
        add_formula(s, Inches(1.4), ft, Inches(10.5), formula)
    set_notes(s, notes)
    return s


def image_slide(num, title, image, bullets=None, caption=None, notes='',
                two=None, base=18):
    s = slide()
    title_band(s, title, num)
    top = Inches(1.35)
    if bullets:
        add_bullets(s, Inches(0.6), top, Inches(12.1), Inches(1.5), bullets, base)
        img_top = Inches(2.95)
    else:
        img_top = Inches(1.5)
    img_h = Inches(7.15) - img_top
    if two:
        # iki görsel yan yana
        add_image_fit(s, image, Inches(0.5), img_top, Inches(6.0), img_h, caption)
        add_image_fit(s, two, Inches(6.9), img_top, Inches(6.0), img_h, None)
    else:
        add_image_fit(s, image, Inches(2.4), img_top, Inches(8.5), img_h, caption)
    set_notes(s, notes)
    return s


# 2 — Akış
content(2, 'Sunum Akışı', [
    'Motivasyon & temel kavramlar (SISO, kapasite)',
    'MIMO ve çok antenin amaçları',
    'Uzamsal çoğullama — kapasite avantajı',
    'V-BLAST mimarisi ve sistem modeli',
    'Alıcı: nulling + sıralı iptal (SIC)',
    'Optimal sıralama (maximin)',
    'Makale laboratuvar sonuçları',
    'Kendi MATLAB simülasyonumuz (SIM-1/2/3)',
    'Sonuç ve günümüze yansıması',
], notes="Yol haritamız şöyle: Önce neden yeni bir yönteme ihtiyaç duyulduğunu, kapasite sorununu "
         "konuşacağız. Ardından MIMO'nun ne olduğunu ve birden fazla anteni hangi amaçlarla "
         "kullanabileceğimizi göreceğiz. Buradan V-BLAST'e geçip alıcının kalbindeki iki fikri — "
         "nulling ve sıralı iptal — adım adım inceleyeceğiz. Makalenin laboratuvar sonuçlarını "
         "gördükten sonra, aynı ortamı MATLAB'de sıfırdan simüle edip kendi sonuçlarımızı sunacağım. "
         "Yaklaşık 40 dakika sürecek; sonunda sorularınız için zaman ayıracağım.")

# 3 — Motivasyon
content(3, 'Motivasyon: Veri Talebi & Spektrum Kıtlığı', [
    'Kablosuz veri talebi katlanarak artıyor',
    'Frekans spektrumu kıt ve pahalı bir kaynak',
    'Daha fazla bant → yok; daha fazla güç → verimsiz',
    '##Asıl soru: Aynı spektrumla daha çok bit nasıl?',
    'V-BLAST\'in cevabı: çok yollu yayılımı (multipath) bir fırsata çevirmek',
], notes="Her yıl daha fazla cihaz, daha fazla veri istiyor. Ama kablosuz haberleşmenin hammaddesi "
         "olan frekans spektrumu kıt bir kaynak; lisansları milyarlarca dolar. İki klasik çözüm var: "
         "ya daha geniş bant kullanırsınız — ama bant yok — ya da daha fazla güç verirsiniz — ama bu "
         "çok verimsiz. Asıl soru şu: aynı frekans bandından, gücü de çok artırmadan, nasıl çok daha "
         "fazla bit geçirebiliriz? V-BLAST'in cevabı, normalde baş belası saydığımız çok yollu "
         "yayılımı bir avantaja çevirmekten geçiyor.")

# 4 — SISO
content(4, 'SISO ve Kablosuz Kanal', [
    'SISO = tek verici, tek alıcı anten',
    'Kanalda: zayıflama + gürültü + çok yollu yayılım (multipath)',
    'Multipath → sönümlenme (fading): sinyal gücü dalgalanır',
    '##SISO\'da multipath bir sorundur — sinyali bozar',
    'Aklınızda tutun: MIMO bu bakışı tersine çevirecek',
], notes="En basit sistemle başlayalım: SISO, yani tek verici ve tek alıcı anten. Sinyal havada yol "
         "alırken zayıflar, gürültüye karışır ve duvarlara çarpıp birden fazla yoldan alıcıya ulaşır — "
         "buna multipath diyoruz. Bu kopyalar bazen birbirini güçlendirir, bazen söndürür; sinyal gücü "
         "dalgalanır, buna sönümlenme yani fading denir. Tek antenli sistemlerde multipath tamamen bir "
         "baş belasıdır. Bu cümleyi aklınızda tutun, çünkü MIMO bunu tersine çevirecek.")

# 5 — Kapasite
content(5, 'Kapasite Sorunu', [
    'Bir kanalın teorik hız sınırı: Shannon kapasitesi',
    'Bant B kıt → artıramıyoruz',
    'SNR logaritmanın içinde → gücü 2× yapmak kapasiteyi az artırır (verimsiz)',
    '##Spektral verimlilik (bps/Hz) geleneksel sistemlerde tıkanır',
], formula=['C  =  B · log₂(1 + SNR)'],
   formula_top=Inches(4.6),
   notes="Bir kanaldan ne kadar hızlı veri geçirebileceğimizin sınırını Shannon formülü verir: kapasite "
         "eşittir bant genişliği çarpı log iki tabanında bir artı SNR. Acı gerçek şu: bandı artıramıyoruz "
         "çünkü spektrum kıt. Gücü iki katına çıkarsak bile SNR logaritmanın içinde olduğu için kapasite "
         "sadece bir tık artar — yani güç pompalamak verimsiz. Sonuçta spektral verimlilik birkaç bps/Hz'de "
         "tıkanır. Peki ne yapacağız? Yepyeni bir boyut eklemek: uzayı.")

# 6 — MIMO nedir
content(6, 'MIMO Nedir?', [
    'MIMO = M verici + N alıcı anten (M ≤ N)',
    'Aynı frekansta, aynı anda, birden çok bağımsız akış',
    '"Karışsınlar — yeter ki alıcıda ayıralım"',
    'Kanal matrisi H\'nin bağımsız sütunları → ayrıştırma mümkün',
    '##Benzetme: yankılı odada 4 konuşmacı, farklı yerlerde 4 mikrofon',
], notes="MIMO, çok girişli çok çıkışlı demek: vericide M, alıcıda N anten. Fikir cüretkâr: aynı "
         "frekansta, aynı anda birden fazla bağımsız akış gönderelim. Sezgi 'karışırlar' der; MIMO "
         "'karışsınlar, alıcıda ayırırız' der. Bunu mümkün kılan kanal matrisi H: bol yansıma varsa her "
         "verici-alıcı çiftinin kanalı farklı olur, H'nin sütunları bağımsızlaşır ve denklem çözülebilir "
         "hale gelir. Yankılı bir odada dört kişi konuşurken farklı yerlerdeki dört mikrofon sesleri "
         "ayırabilir — MIMO tam olarak budur.")

# 7 — Çok antenin üç amacı
s = slide()
title_band(s, 'Çok Antenin ÜÇ Amacı', 7)
add_bullets(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.0), [
    '##1) Çeşitleme (Diversity) — GÜVENİLİRLİK için',
    ('Aynı veriyi farklı bağımsız yollardan gönder/al → sönümlenmeye dayanıklılık', 1),
    ('RX diversity (MRC) ve TX diversity (Alamouti)', 1),
    '##2) Hüzmeleme (Beamforming) — SNR / KAPSAMA için',
    ('Enerjiyi belirli yöne odakla; istenmeyen yöne null koy', 1),
    '##3) Uzamsal Çoğullama (SM) — HIZ için',
    ('Farklı antenden farklı veri → hız anten sayısıyla katlanır', 1),
    ('►  V-BLAST tam olarak BUDUR', 1),
    '##Ödünleşim: aynı antenler ya hıza ya güvenilirliğe (V-BLAST hıza odaklanır, N>M ile biraz çeşitleme de kazanır)',
], base=18)
set_notes(s, "Önemli bir noktayı netleştirelim: birden fazla anten kullanmak tek bir şey değildir; üç amacı "
             "vardır. Birincisi çeşitleme, diversity: aynı veriyi farklı yollardan gönderip alırsınız, biri "
             "sönümlenmeye düşse bile diğeri sağlam kalır — güvenilirlik artar. Alıcıda yaparsanız RX diversity, "
             "vericide Alamouti gibi kodlarla yaparsanız TX diversity. İkincisi hüzmeleme: enerjiyi bir yöne "
             "odaklar, SNR ve kapsama artar. Üçüncüsü uzamsal çoğullama: farklı antenden farklı veri gönderip "
             "hızı katlarsınız — V-BLAST tam olarak budur. Ama aynı antenleri ya hıza ya güvenilirliğe "
             "ayırırsınız; bir ödünleşim var. V-BLAST hıza odaklanır, ama alıcı anten fazlalığıyla biraz "
             "çeşitleme de kazanır.")

# 8 — Uzamsal çoğullama
content(8, 'Uzamsal Çoğullama & Kapasite Avantajı', [
    '1 akışı M alt-akışa böl, ayrı antenlerden aynı anda gönder',
    'Kapasite anten sayısı M ile DOĞRUSAL büyür',
    'Ekstra bant veya ekstra güç GEREKMEZ',
    '##V-BLAST = bu fikrin gerçek donanımda çalışan hali',
], formula=['C_MIMO  ≈  M · B · log₂(1 + SNR)',
            '1 anten → 1× ,   8 anten → ~8× kapasite'],
   formula_top=Inches(4.7),
   notes="Çoğullamanın gücü şurada: SISO'da gücü iki katına çıkarmak kapasiteyi sadece logaritmik "
         "artırıyordu. MIMO'da ise kapasite anten sayısı M ile doğrudan çarpılıyor — doğrusal büyüyor. "
         "Sekiz anten kabaca sekiz kat kapasite, üstelik ne ekstra bant ne ekstra güç gerekmeden. "
         "Spektrumun altın değerinde olduğu bir dünyada bu bir devrim. V-BLAST bu fikrin gerçek "
         "donanımda çalışan pratik gerçeklemesi. Peki bu karışmış akışları alıcıda nasıl ayıracağız?")

# 9 — BLAST D vs V
content(9, 'BLAST Ailesi: D-BLAST vs V-BLAST', [
    '##D-BLAST (Diagonal)',
    ('Alt-akışlar arası çapraz kodlama; ~%90 Shannon', 1),
    ('Karmaşık → gerçek zamanlı uygulaması zor', 1),
    '##V-BLAST (Vertical)',
    ('Kodlama YOK — sadece "böl ve gönder"', 1),
    ('Biraz düşük ama yüksek verim; GERÇEK ZAMANLI uygulandı', 1),
    '##V-BLAST\'in çözdüğü problem: MIMO verimini PRATİK, çalışan bir alıcıyla elde etmek',
    'Verici basit; tüm zekâ alıcıda',
], base=19,
   notes="BLAST, Bell Labs Layered Space-Time'ın kısaltması ve iki üyesi var. D-BLAST, alt-akışlar "
         "arasında çapraz kodlama kullanır; teorik olarak müthiş, Shannon'un yüzde doksanına ulaşır, ama "
         "gerçek zamanlı uygulaması çok zor. V-BLAST bu kodlamayı atar: sadece böl ve gönder. Biraz "
         "performanstan feragat eder ama gerçek zamanlı çalışan basit bir sistem sunar. Çözdüğü problem "
         "net: MIMO verimini pratik, çalışan bir alıcıyla elde etmek. Vericide her şey basit; tüm zekâ "
         "alıcıda.")

# 10 — Sistem mimarisi (Şekil 1)
image_slide(10, 'V-BLAST Sistem Mimarisi', F1, bullets=[
    'Vektör kodlayıcı → M verici (her biri sıradan QAM) → zengin saçılma → N alıcı → V-BLAST işleme',
    'Toplam güç sabit: her verici 1/M güçle yayın yapar',
], caption='Şekil 1 — V-BLAST yüksek seviye sistem diyagramı (makaleden)',
   notes="Sistemin kuş bakışı görünümü bu. Tek veri akışı vektör kodlayıcıda M alt-akışa bölünüyor, her "
         "biri sıradan bir QAM vericisine gidiyor. M verici aynı frekansta, senkronize, aynı anda yayın "
         "yapıyor; toplam güç sabit kalsın diye her biri 1/M güçle. Sinyaller zengin saçılma ortamından "
         "geçip karışarak N alıcı antene ulaşıyor. Bütün marifet sağdaki kutuda: V-BLAST sinyal işleme. "
         "Sunumun geri kalanı bu kutunun içini açmak olacak.")

# 11 — Sistem modeli
content(11, 'Sistem Modeli & Varsayımlar', [
    'a: gönderilen semboller (M×1) — bulmak istediğimiz',
    'H: kanal matrisi (N×M) — eğitim dizisiyle tahmin edilir',
    'r: alınan vektör (N×1) ,  ν: gürültü (varyans σ²)',
    '##Flat fading, yarı-durağan kanal (burst boyunca sabit)',
    'Problem: r ve H biliniyor → a\'yı çöz (gürültü yüzünden basit değil)',
], formula=['r  =  H · a  +  ν'],
   formula_top=Inches(5.2),
   notes="Tüm problemi tek denkleme indirgeyelim: alınan vektör r eşittir kanal H çarpı gönderilen a, "
         "artı gürültü nü. Kanal bir burst boyunca sabit kalıyor — yarı-durağan — ve eğitim dizisiyle H'yi "
         "doğru tahmin ediyoruz. Yani r ve H biliniyor, a'yı çözeceğiz. 'a eşittir H ters çarpı r' demek "
         "istiyorsunuz; doğru ama gürültü yüzünden iş o kadar basit değil — alıcı tasarımı burada devreye "
         "giriyor.")

# 12 — Nulling
content(12, 'Detection & Nulling (ZF / MMSE)', [
    'Her akışı sırayla "istenen", gerisini "girişim" say; w ile bastır',
    'ZF (Zero-Forcing): diğerlerini tam sıfırlar — basit ama gürültüyü büyütür',
    'MMSE: gürültü + girişimi birlikte dengeler — genelde daha iyi',
    '##Karar istatistiği:  y_i = wᵢᵀ · r',
], formula=['ZF koşulu:   wᵢᵀ · (H)ⱼ  =  δᵢⱼ      (i=j → 1 ,  i≠j → 0)'],
   formula_top=Inches(5.1),
   notes="Karışmış sinyalleri ayırmanın ilk yolu nulling. Bir akışı istenen, diğerlerini girişim kabul "
         "ediyoruz ve alınan sinyali öyle bir ağırlık vektörüyle çarpıyoruz ki istenen geçsin, diğerleri "
         "sıfırlansın. Zero-Forcing'de ağırlık vektörü istenen sütunla çarpınca bir, diğerleriyle sıfır "
         "verir — formüldeki Kronecker delta budur. ZF basit ama gürültüyü büyütebilir. MMSE ise girişim "
         "ile gürültüyü birlikte dengeler, genelde daha iyidir; bunu simülasyonda karşılaştıracağız. "
         "Ama nulling tek başına en iyisi değil.")

# 13 — SIC
content(13, 'Sembol İptali (SIC): Adım Adım', [
    '##Döngü (her akış için):',
    ('1) Nulling → karar istatistiği y', 1),
    ('2) Slice:  â = Q(y)  (en yakın QAM noktası)', 1),
    ('3) Cancel: çözülen sembolü sinyalden çıkar', 1),
    ('4) Kalan akışlar için tekrarla', 1),
    '##Her turda bir girişim daha az → problem küçülür (DFE\'ye benzer)',
], formula=['r_{i+1}  =  r_i  −  â_{k_i} · (H)_{k_i}'],
   formula_top=Inches(5.2),
   notes="V-BLAST'i güçlü kılan ikinci fikir: ardışık sembol iptali, SIC. Bir sembolü çözdüysem, onun "
         "sinyale kattığı katkıyı çıkarabilirim; geriye bir girişim daha az kalır. Döngü: nulling ile kaba "
         "tahmin al, slice et yani en yakın noktaya yuvarla, sonra bu tahmini kanal sütunuyla çarpıp "
         "sinyalden çıkar. Problem küçüldü, kalan akışlar için tekrarla. Her turda iş kolaylaşıyor. Bu, "
         "karar geri-beslemeli denkleştirmeye benzer. Peki neden saf nulling'den iyi?")

# 14 — Neden iptal daha iyi
content(14, 'Neden İptal Daha İyi?', [
    'Saf nulling: her w, M−1 akışa ortogonal olmalı (ağır kısıt)',
    'İptal ile: w yalnızca KALAN akışlara ortogonal (az kısıt)',
    'Cauchy-Schwarz: çok kısıt → büyük ‖w‖',
    '##Küçük norm = yüksek SNR → daha az hata',
], formula=['ρ_{k_i}  =  ⟨|a_{k_i}|²⟩ / ( σ² · ‖w_{k_i}‖² )'],
   formula_top=Inches(5.1),
   notes="Cevap ağırlık vektörünün normunda. Saf nulling'de her vektör diğer M eksi bir akışın hepsine "
         "ortogonal olmak zorunda — ağır kısıt. İptalle, bir akışı çıkardıktan sonra vektör yalnızca kalan "
         "akışlara ortogonal olmak zorunda — az kısıt. Cauchy-Schwarz gereği, bir vektör ne kadar çok şeye "
         "ortogonal olmaya zorlanırsa normu o kadar büyür. Ve şu formül kritik: detection sonrası SNR, "
         "normun karesiyle ters orantılı. Küçük norm, yüksek SNR. İptal kısıtları azaltır, norm küçülür, "
         "SNR yükselir, hata düşer.")

# 15 — Optimal sıralama
content(15, 'Optimal Sıralama: "En İyiyi Önce Al"', [
    'İptalde SIRA önemli (saf nulling\'de değil)',
    'En zayıf akış hatayı belirler → hedef: minimumu maksimize et (maximin)',
    '##Şaşırtıcı sonuç: her adımda en güçlü akışı seçmek → KÜRESEL optimum',
    'En güvenilir kararı başta ver → hata yayılımı azalır',
    'Optimalliği ilk kez bu makalede ispatlandı',
], formula=['k_i  =  arg min_j  ‖ (G_i)_j ‖²     (en küçük normlu satır)'],
   formula_top=Inches(5.3),
   notes="İptal kullanınca yeni soru: hangi sırayla çözelim? Önce çözülen akış sonrakileri etkiler. En "
         "zayıf akış hatayı belirlediği için en kötü akışın SNR'ını maksimize etmek isteriz — maximin. "
         "Makalenin zarif sonucu: her adımda o anki en güçlü akışı seçmek, küresel olarak en iyi sıralamayı "
         "veriyor. Pratikte en güçlü akış, pseudoinverse'in en küçük normlu satırına karşılık gelir. Bir "
         "faydası daha: en güvenilir kararı başta verince hata yayılımını da azaltıyoruz. Bu sezgisel kural "
         "yıllardır kullanılıyordu ama optimalliği ilk kez bu makalede ispatlandı.")

# 16 — Tam algoritma
s = slide()
title_band(s, 'Tam Algoritma (Özyinelemeli) — Eq. (9)', 16)
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.0), [
    '##Başlangıç: G₁ = H⁺ ,  ilk akış = en küçük normlu satır',
], base=19)
add_formula(s, Inches(1.0), Inches(2.4), Inches(11.3), [
    'w_{k_i} = (G_i)_{k_i}            (nulling vektörü)',
    'y_{k_i} = w_{k_i}ᵀ · r_i          (karar istatistiği)',
    'â_{k_i} = Q( y_{k_i} )            (slicing)',
    'r_{i+1} = r_i − â_{k_i}·(H)_{k_i}    (iptal / cancellation)',
    'G_{i+1} = (H_{k̄_i})⁺             (deflation: sütunu sıfırla)',
    'k_{i+1} = arg min_j ‖(G_{i+1})_j‖²   (sonraki en iyi akış)',
], size=17)
set_notes(s, "Bütün parçaları tek özyinelemeli algoritmada birleştirelim. Başlangıçta H'nin "
             "pseudoinverse'ini hesaplayıp en küçük normlu satırdan ilk akışı seçiyoruz. Döngü: nulling "
             "vektörünü al, karar istatistiğini hesapla, slice et, tahmini çıkar. Önemli adım: H'yi "
             "sönükleştiriyoruz — çözdüğümüz sütunu sıfırlayıp pseudoinverse'i yeniden hesaplıyoruz, çünkü "
             "o akış artık yok. Yeni matriste yine en küçük normlu satırı seçip devam ediyoruz. Bütün "
             "akışlar bitene kadar. Birkaç satır gibi görünse de bu, MIMO alıcısının kalbidir. Şimdi "
             "gerçekten çalışıyor mu, bakalım.")

# 17 — Makale sonuçları (Şekil 2 + 3)
image_slide(17, 'Makale Laboratuvar Sonuçları', F2, two=F3, bullets=[
    'Prototip: 1.9 GHz, M=8/N=12, 16-QAM, gerçek ofis ortamı',
    'Sol (Şekil 2): iptal ~4 dB kazanç • Sağ (Şekil 3): konuma göre gürbüz • 20–40 bps/Hz, 621 kbps@30kHz',
], caption='Şekil 2 — BLER vs SNR',
   notes="Makale sadece teori değil; gerçek bir prototip kurulmuş. 1.9 gigahertz'de, sekiz verici on iki "
         "alıcı antenle, gerçek ofis ortamında. Soldaki Şekil 2'de üstteki eğri saf nulling, alttaki "
         "nulling artı optimal sıralı iptal; aradaki fark yaklaşık dört desibel — bu konfigürasyonda "
         "neredeyse on bps/Hz'lik verimlilik farkı. Sağdaki Şekil 3, vericiyi farklı konumlara "
         "taşıdıklarında performansın gürbüz kaldığını gösteriyor. Sonuç: 20 ila 40 bps/Hz, sadece 30 "
         "kilohertz bantta 621 kilobit. Aynı verimliliği tek antenle elde etmek için milyarlarca noktalı "
         "bir takımyıldız gerekirdi. Şimdi bunu kendimiz doğrulayalım.")

# 18 — Simülasyon kurulumu
content(18, 'Bizim Simülasyonumuz: Kurulum', [
    'Makale ortamı MATLAB\'de — TOOLBOX FONKSİYONU KULLANMADAN',
    'M=8, N=12, 16-QAM, Rayleigh flat fading (H ~ i.i.d. CN(0,1))',
    'Sıfırdan: QAM mod/demod, kanal, gürültü, ZF/MMSE/SIC alıcı',
    'Monte Carlo: BLER & BER vs SNR',
    '##Göstereceğimiz: SIM-1 ana karşılaştırma · SIM-2 sıralama · SIM-3 alıcı tipi + N etkisi',
], base=19,
   notes="Makaledeki ortamı doğrulamak için her şeyi MATLAB'de sıfırdan, hiçbir hazır toolbox fonksiyonu "
         "kullanmadan kodladık: 16-QAM modülasyonu, kanal, gürültü, pseudoinverse, sıralı iptalli alıcı — "
         "hepsi elle. Parametreler makaleyle aynı: sekiz verici, on iki alıcı, 16-QAM, bağımsız Rayleigh "
         "kanal. Her SNR noktasında binlerce burst üretip Monte Carlo ile blok ve bit hata oranını ölçtük. "
         "Sırasıyla dört şey göstereceğim: ana karşılaştırma, sıralamanın etkisi, alıcı tipinin etkisi ve "
         "anten sayısının etkisi.")

# 19 — SIM-1
image_slide(19, 'Simülasyon 1: BLER & BER — Nulling vs V-BLAST', SIM1, bullets=[
    'V-BLAST, saf nulling\'den ~4 dB daha iyi (ölçülen: 4.4 dB) → makale ~4 dB ile birebir',
    'V-BLAST eğimi daha dik: her doğru karar sonrakini kolaylaştırır',
], caption='SIM-1 — kendi kodumuz, M=8 N=12 16-QAM',
   notes="İşte ana sonucumuz, makaledeki Şekil 2'nin muadili. Yatay eksende SNR, dikey eksende hata oranı, "
         "logaritmik ölçekte. Üstteki saf nulling, alttaki V-BLAST. V-BLAST eğrisi yaklaşık dört desibel "
         "solda; ölçtüğümüz kazanç 4.4 desibel — makalenin yaklaşık dört desibeliyle neredeyse birebir aynı. "
         "V-BLAST'in eğimi de daha dik, çünkü her doğru kararla bir girişimi temizleyip sonrakini "
         "kolaylaştırıyor. Makalenin ana iddiasını kendi kodumuzla, sıfırdan doğruladık.")

# 20 — SIM-2
image_slide(20, 'Simülasyon 2: Sıralamanın Etkisi', SIM2, bullets=[
    'Sıralı (maximin) iptal, sırasıza göre ~3 dB daha iyi',
    'Makalenin ana teorik katkısının (optimal sıralama) deneysel kanıtı',
], caption='SIM-2 — sıralı SIC vs sabit sıralı SIC',
   notes="İkinci grafik makalenin en zarif sonucunu test ediyor: sıralama önemli mi? İki iptal "
         "stratejisini karşılaştırıyoruz. Üstteki eğri akışları sabit, gelişigüzel sırayla çözüyor; "
         "alttaki V-BLAST'in optimal sıralaması: her adımda en güçlü akışı önce çöz. Aradaki fark yaklaşık "
         "üç desibel. Bu, 'en iyiyi önce al' kuralının değerinin doğrudan kanıtı: en güvenilir kararlar "
         "başta verilince hata yayılımı azalıyor, en zayıf akışın SNR'ı yükseliyor. Makalenin matematiksel "
         "olarak ispatladığı maximin optimalliğini deneysel olarak görüyoruz.")

# 21 — SIM-3
image_slide(21, 'Simülasyon 3: Alıcı Tipi & Anten Sayısı', SIM3, bullets=[
    '(a) N=M=8 kare sistemde MMSE, ZF\'ten ~3 dB iyi (ZF gürültüyü büyütür)',
    '(b) N=12→16→20 arttıkça performans belirgin iyileşir (çeşitleme kazancı)',
], caption='SIM-3 — (a) ZF vs MMSE   (b) N etkisi',
   notes="Son simülasyonda iki tasarım sorusu. Solda zorlu bir durum seçtim: alıcı anten sayısı verici "
         "sayısına eşit, N eşittir M eşittir sekiz, kare sistem. Bu durumda kanal kötü koşullu olur ve ZF "
         "gürültüyü ciddi büyütür; MMSE bit hata oranında ZF'ten yaklaşık üç desibel daha iyi, çünkü "
         "girişim ile gürültüyü birlikte dengeler. Sağda alıcı anten sayısının etkisi: M sabitken N'i on "
         "ikiden yirmiye çıkardığımızda performans belirgin iyileşiyor. Sebep çeşitleme: her fazladan "
         "alıcı anten ekstra serbestlik ve sönümlenmeye dayanıklılık getiriyor. Bu da V-BLAST'in neden "
         "N'i M'den büyük tuttuğunu açıklıyor.")

# 22 — Sonuç
content(22, 'Sonuçlar, Önem & Kapanış', [
    'V-BLAST: basit verici + akıllı alıcı (nulling + sıralı iptal + optimal sıra)',
    '20–40 bps/Hz — o güne dek eşi görülmemiş; multipath\'i fırsata çevirdi',
    'Simülasyonumuz doğruladı: V-BLAST>nulling, sıralı>sırasız, MMSE>ZF, N↑',
    '##Bugün: Wi-Fi (802.11n/ac/ax), 4G LTE ve 5G MIMO\'sunun temeli',
    'Teşekkürler — sorularınız?',
], base=20,
   notes="Toparlayalım. V-BLAST'in dehası, vericiyi basit tutup tüm zekâyı alıcıya yüklemesinde: nulling "
         "ile ayır, sıralı iptal ile temizle, her zaman en güçlü akışı önce çöz. Böylece düşmanımız olan "
         "çok yollu yayılımı avantaja çevirip 20 ila 40 bps/Hz gibi görülmemiş verimliliklere ulaşıldı. "
         "Kendi sıfırdan yazdığımız MATLAB simülasyonu da bütün ana sonuçları doğruladı: V-BLAST nulling'i, "
         "sıralı sırasızı, MMSE ZF'i yendi ve anten sayısı arttıkça performans iyileşti. Ve bu akademik bir "
         "merak değil: bugünkü Wi-Fi, 4G ve 5G MIMO'sunun doğrudan temeli bu çalışma. Teşekkürler, "
         "sorularınızı alabilirim.")

# 23 — Kaynaklar
content(23, 'Kaynaklar & Sorular', [
    'Wolniansky, Foschini, Golden, Valenzuela — V-BLAST (1998)',
    'G. J. Foschini — Layered Space-Time Architecture, Bell Labs Tech. J. (1996)',
    'Foschini & Gans — On Limits of Wireless Comm. (1998)',
    '##Sorular?',
    'Hazır Q&A: D-BLAST farkı · ZF/MMSE · sıralama · hata yayılımı · zengin saçılma · M≤N · çok antenin amaçları',
], base=19,
   notes="Ana kaynaklar bunlar; özellikle Foschini'nin 1996 temel makalesi ve elimizdeki V-BLAST "
         "çalışması. Sorularınız için buradayım.")

out = os.path.join(HERE, 'V-BLAST_Sunum.pptx')
prs.save(out)
print('[+] Kaydedildi:', out, '| slayt sayısı:', len(prs.slides._sldIdLst))
